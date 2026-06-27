"""Stage 3 — Deck creation (post Phase 1 integration).

Architecture: Python-orchestrated, locally executed. NO sandbox, NO base64
extraction, NO reference-deck uploads at runtime.

For each slide in the spec:
  - Templated types (title, section dividers, company context, next steps)
    → rendered directly by ``templates.py``, no LLM call.
  - Non-templated types (opportunity buckets, deep-dives, engagement stages,
    workshop, commercial close) → Sonnet generates python-pptx code for
    THAT ONE SLIDE; we ``exec`` it locally on the in-memory ``Presentation``.
    On code error, one retry with the traceback included.

Per-slide cost is bounded; total Stage 3 cost is bounded; the runaway iteration
loops that produced $20 runs are no longer possible.
"""

from __future__ import annotations

import json
import os
import re
import traceback
import zipfile
from typing import Callable

import anthropic
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

import prompts
import templates
from api_retry import with_retries
from usage_tracker import StageBudgetExceeded, check_budget

# Sonnet for per-slide code generation. Opus is overkill for single-slide layout
# work and would defeat the cost-reduction purpose of this architecture.
SLIDE_MODEL = "claude-sonnet-4-6"

STAGE_3_BUDGET_USD = float(os.environ.get("STAGE_3_BUDGET_USD", "4.0"))
MAX_SLIDE_RENDER_ATTEMPTS = 2  # initial + 1 retry on error

_SECTION_KEYS = ("section_1", "section_2", "section_3", "section_4")
_SECTION_NAMES = ("Company context", "Opportunity buckets", "Deep-dives", "Engagement")
_SECTION_NUMBERS = ("01", "02", "03", "04")


_CITATION_RE = re.compile(
    r"\s*\[\s*(?:Sources?:\s*[^\]]*|S\d+(?:\s*[,;]\s*S\d+)*)\s*\]\s*",
    re.IGNORECASE,
)


def _strip_citations(text: str | None) -> str:
    """Remove inline '[Sources: S3, S7]' style citations from slide body text."""
    if not text:
        return ""
    return _CITATION_RE.sub(" ", str(text)).strip()


# Catalogue-internal keys that must never reach a rendered slide. The slide renderer
# only needs display fields (name, headline, prose, roi, etc.); these are organising
# metadata. Belt-and-suspenders behind the prompt rule — the renderer can't print what
# it never receives.
_INTERNAL_KEYS = {"id", "opportunity_id", "sub_pillar", "form_pillar"}


def _scrub_spec(spec: dict) -> dict:
    """Recursively strip [Sources: ...] markers AND internal catalogue ids/pillar keys
    from a slide spec, so neither can leak onto a rendered slide."""
    if isinstance(spec, dict):
        return {k: _scrub_spec(v) for k, v in spec.items() if k not in _INTERNAL_KEYS}
    if isinstance(spec, list):
        return [_scrub_spec(v) for v in spec]
    if isinstance(spec, str):
        return _strip_citations(spec)
    return spec


# --------------------------------------------------------------------------------- #
# Per-slide LLM rendering — Sonnet returns python-pptx code; we exec it
# --------------------------------------------------------------------------------- #
_CODE_FENCE_RE = re.compile(
    r"`{3,}[ \t]*\w*[ \t]*\r?\n([\s\S]*?)`{3,}",
    re.MULTILINE,
)


def _extract_code(text: str) -> str:
    """Pull a python code block out of the model response. Tolerant of fence variations:
    3+ backticks, optional language tag (python/py/etc.), CRLF or LF endings, leading
    whitespace, and missing closing fence."""
    text = (text or "").strip()

    # Primary: properly fenced block.
    m = _CODE_FENCE_RE.search(text)
    if m:
        return m.group(1).strip()

    # Fallback: strip any leading or trailing fence lines manually. Handles the case
    # where the model opens a fence but never closes it (truncated output), or where
    # the fence is on its own line surrounded by extra whitespace.
    lines = text.splitlines()
    while lines and lines[0].lstrip().startswith("```"):
        lines.pop(0)
    while lines and lines[-1].rstrip().endswith("```"):
        stripped = lines[-1].rstrip().rstrip("`").rstrip()
        if stripped:
            lines[-1] = stripped
            break
        lines.pop()
    return "\n".join(lines).strip()


def _text_of(message) -> str:
    return "".join(b.text for b in message.content if getattr(b, "type", None) == "text")


def _render_slide_via_claude(
    client: anthropic.Anthropic,
    prs: Presentation,
    slide_spec: dict,
    brand: dict,
    section_key: str,
    slide_type: str,
    tracker=None,
    say: Callable[[str], None] = lambda _m: None,
) -> None:
    """Get python-pptx code from Sonnet for one slide; exec it locally on ``prs``."""
    user_prompt = prompts.slide_render_prompt(slide_spec, brand, section_key, slide_type)

    messages = [{"role": "user", "content": user_prompt}]
    last_error: str | None = None

    for attempt in range(MAX_SLIDE_RENDER_ATTEMPTS):
        resp = with_retries(lambda: client.messages.create(
            model=SLIDE_MODEL,
            max_tokens=4000,
            system=prompts.SLIDE_RENDER_SYSTEM_PROMPT,
            messages=messages,
        ))

        # Record usage and enforce the cap on EVERY call (not just end of slide).
        if tracker is not None and getattr(resp, "usage", None) is not None:
            tracker.record("stage_3", SLIDE_MODEL, resp.usage)
            check_budget(tracker, "stage_3", STAGE_3_BUDGET_USD)

        if resp.stop_reason == "refusal":
            raise RuntimeError(f"Stage 3 slide render ({slide_type}) was refused.")

        code = _extract_code(_text_of(resp))
        if not code:
            last_error = "Model returned no code block."
            messages.append({"role": "assistant", "content": resp.content})
            messages.append({"role": "user", "content":
                "You did not return a python code block. Return ONLY one "
                "```python ... ``` code block that adds the slide to `prs`. "
                "No prose."})
            continue

        # Snapshot slide count so we can roll back a partial slide on error.
        before = len(prs.slides)

        try:
            exec_globals = {
                "prs": prs,
                "brand": brand,
                "slide_spec": slide_spec,
                "Presentation": Presentation,
                "Inches": Inches,
                "Pt": Pt,
                "Emu": Emu,
                "RGBColor": RGBColor,
                "PP_ALIGN": PP_ALIGN,
                "MSO_ANCHOR": MSO_ANCHOR,
                "MSO_SHAPE": MSO_SHAPE,
            }
            exec(code, exec_globals)

            # Sanity: the call should have added exactly one slide.
            added = len(prs.slides) - before
            if added != 1:
                raise RuntimeError(
                    f"Slide code added {added} slides, expected 1. "
                    f"Code must add exactly one slide to `prs`."
                )
            return  # success
        except Exception as exc:
            last_error = traceback.format_exc()
            # Roll back any partially-added slides so the deck stays clean.
            while len(prs.slides) > before:
                _remove_last_slide(prs)
            messages.append({"role": "assistant", "content": resp.content})
            messages.append({"role": "user", "content":
                f"Your code errored when I executed it:\n\n```\n{last_error}\n```\n\n"
                "Fix it and return ONE python code block only. The code must add "
                "EXACTLY one slide to `prs`."})

    raise RuntimeError(
        f"Stage 3: per-slide rendering failed for {slide_type} after "
        f"{MAX_SLIDE_RENDER_ATTEMPTS} attempts.\n\nLast error:\n{last_error}"
    )


def _remove_last_slide(prs: Presentation) -> None:
    """Remove the last slide. python-pptx has no public remove, so we hit XML."""
    slides_xml = prs.slides._sldIdLst
    if len(slides_xml) == 0:
        return
    last = slides_xml[-1]
    rId = last.attrib.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    slides_xml.remove(last)
    if rId:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


# --------------------------------------------------------------------------------- #
# Slide-type classification
# --------------------------------------------------------------------------------- #
def _classify(section_key: str, slide_spec: dict) -> str:
    """Decide which slide_type the renderer should treat this as.

    Section 1 → 'company_context' (templated).
    Section 2 → 'opportunity_bucket'.
    Section 3 → 'deep_dive'.
    Section 4 → 'next_steps' — the only Section 4 slide type. Engagement-model,
                workshop, and commercial-close layouts were retired; Section 4 is a
                single action-oriented Next Steps close (see _select_section_4_slide).
    """
    if section_key == "section_1":
        return "company_context"
    if section_key == "section_2":
        return "opportunity_bucket"
    if section_key == "section_3":
        return "deep_dive"
    return "next_steps"


def _synthesize_next_steps(company_name: str) -> dict:
    """Generic Next Steps spec used when Stage 2 didn't produce one.

    Action-oriented close: three concrete asks of the client, two commitments from
    Aistra. Generic enough to apply across companies; the user can edit post-gen.
    """
    return {
        "slide_type": "next_steps",
        "title": "Next steps",
        "subtitle": f"Three with {company_name}, two with Aistra.",
        "with_company": [
            f"Identify the internal sponsor and finance counterpart for Discovery scoping.",
            f"Confirm the data access pathway for the analyses scoped in this deck.",
            f"Decide the Discovery kickoff window — six weeks from green-light.",
        ],
        "with_aistra": [
            "Circulate the Discovery scope and data requirements list within five business days.",
            "Share the Prototype design for the lead opportunity at kickoff.",
        ],
    }


def _has_next_steps_shape(s: dict) -> bool:
    """True if a Stage-2 slide carries Next Steps content."""
    if not isinstance(s, dict):
        return False
    if s.get("with_company") or s.get("with_aistra") or \
       s.get("client_actions") or s.get("aistra_actions"):
        return True
    stype = (s.get("slide_type") or s.get("type") or "").lower()
    return "next" in stype and "step" in stype


def _select_section_4_slide(slides: list, company_name: str = "") -> list:
    """Section 4 collapses to exactly one slide — always a Next Steps close.

    Use the first Stage-2 slide that carries Next Steps content (with_company /
    with_aistra, or a next_steps slide_type); otherwise synthesize a generic one.
    Any other Section-4 shapes Stage 2 might emit are dropped by design — the deck
    closes on what happens after the meeting, not an engagement-model recap.
    """
    for s in slides or []:
        if _has_next_steps_shape(s):
            return [s]
    return [_synthesize_next_steps(company_name)]


# --------------------------------------------------------------------------------- #
# Main entry point
# --------------------------------------------------------------------------------- #
def run(
    deck_spec: dict,
    brand: dict,
    deck_path: str,
    client: anthropic.Anthropic,
    progress: Callable[[str], None] | None = None,
    tracker=None,
) -> str:
    """Assemble the deck locally; LLM is only used per-slide for non-templated types."""
    say = progress or (lambda _msg: None)

    # Set up the Presentation with 16:9 dimensions from the brand JSON.
    sw = float(brand.get("slide_width_inches", 13.333))
    sh = float(brand.get("slide_height_inches", 7.5))
    prs = Presentation()
    prs.slide_width = Inches(sw)
    prs.slide_height = Inches(sh)

    title_content = deck_spec.get("title_slide", {}) or {}
    thesis = deck_spec.get("thesis", {}) or {}
    sections = deck_spec.get("sections", {}) or {}
    company_name = title_content.get("title", "Company")

    # Footer bookkeeping — each entry maps a slide index to a section index.
    slide_section_index: list[int] = []

    # --- Slide 1: Title (templated) -------------------------------------------- #
    say("Assembling deck — title slide…")
    templates.render_title_slide(
        prs, title_content, thesis, brand,
        audience_role=_audience_label(deck_spec),
        date_str=_today_str(deck_spec),
    )
    slide_section_index.append(-1)  # title slide: no section label in the footer

    # --- Iterate sections ------------------------------------------------------ #
    # Per-section storyline arguments (from thesis_v1_step_b) drive each divider's
    # transition line. Ordered Section 1..4; guard against a short/missing list.
    storyline_args = [
        (item.get("argument") or "").strip()
        for item in (thesis.get("storyline") or [])
        if isinstance(item, dict)
    ]

    for section_idx, section_key in enumerate(_SECTION_KEYS):
        section = sections.get(section_key, {}) or {}
        section_num = _SECTION_NUMBERS[section_idx]
        section_label = _SECTION_NAMES[section_idx]

        # Section divider (templated). Prefer the storyline argument for this section;
        # fall back to any intro/subtitle the section carried, then to blank.
        say(f"Assembling deck — section {section_num} divider…")
        transition = storyline_args[section_idx] if section_idx < len(storyline_args) else ""
        transition = transition or section.get("intro") or section.get("subtitle") or ""
        templates.render_section_divider(prs, section_num, section_label, transition, brand)
        slide_section_index.append(section_idx)

        # Content slides — Section 4 collapses to a single Next Steps slide
        section_slides = section.get("slides") or []
        if section_key == "section_4":
            section_slides = _select_section_4_slide(section_slides, company_name)
        for i, raw_slide_spec in enumerate(section_slides, start=1):
            slide_spec = _scrub_spec(raw_slide_spec)  # remove [Sources: ...] artifacts
            slide_type = _classify(section_key, slide_spec)
            say(f"Assembling deck — {section_label.lower()} slide {i} ({slide_type})…")

            if slide_type == "company_context":
                templates.render_company_context(prs, slide_spec, brand)
            elif slide_type == "next_steps":
                templates.render_next_steps(prs, slide_spec, brand, company_name=company_name)
            elif slide_type == "deep_dive":
                templates.render_deep_dive(prs, slide_spec, brand,
                                           index=i, total=len(section_slides))
            else:
                # Only opportunity_bucket reaches Sonnet now (Sonnet, bounded cost, local exec)
                _render_slide_via_claude(
                    client, prs, slide_spec, brand,
                    section_key=section_key,
                    slide_type=slide_type,
                    tracker=tracker,
                    say=say,
                )

            slide_section_index.append(section_idx)

    # --- Uniform footer + page numbers (added last so totals are right) --------- #
    say("Applying uniform footer and page numbers…")
    templates.apply_uniform_footer(
        prs, company_name, list(_SECTION_NAMES), slide_section_index, brand,
    )

    # --- Save locally ---------------------------------------------------------- #
    os.makedirs(os.path.dirname(deck_path) or ".", exist_ok=True)
    prs.save(deck_path)

    if not zipfile.is_zipfile(deck_path):
        raise RuntimeError("Saved deck is not a valid .pptx (OOXML zip).")

    say(f"Deck assembled and saved ({len(prs.slides)} slides).")
    return deck_path


# --------------------------------------------------------------------------------- #
# Small helpers
# --------------------------------------------------------------------------------- #
def _audience_label(deck_spec: dict) -> str:
    """Best-effort recovery of audience role for the title slide footer line."""
    ts = deck_spec.get("title_slide", {}) or {}
    return ts.get("audience_role") or ts.get("role") or ""


def _today_str(deck_spec: dict) -> str:
    ts = deck_spec.get("title_slide", {}) or {}
    return ts.get("date") or ""


# =================================================================================== #
# v2 — assemble a deck directly from the Storyline (Contract 5). No section dividers;
# deterministic dispatch to the house-style renderers in templates.py.
# =================================================================================== #
def run_v2(storyline: dict, brand: dict, deck_path: str, progress=None) -> dict:
    import templates as T
    from pptx import Presentation
    from pptx.util import Inches

    say = progress or (lambda _m: None)
    say("Assembling deck…")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    meta = storyline.get("meta", {})
    T.render_title_slide(prs, storyline.get("title", {"title": meta.get("company", "")}), meta, brand,
                         audience_role=meta.get("audience_role", "") or "",
                         date_str=meta.get("date", "") or "")

    sc = storyline.get("strategic_context") or {}
    if sc:
        if sc.get("layout") == "compare":
            T.render_strategic_context_compare(prs, sc, brand)
        else:
            T.render_strategic_context_cards(prs, sc, brand)

    if storyline.get("art_of_the_possible"):
        T.render_art_of_possible(prs, storyline["art_of_the_possible"], brand)
    if storyline.get("opportunity_areas"):
        T.render_opportunity_areas(prs, storyline["opportunity_areas"], brand)

    dds = storyline.get("deep_dives") or []
    roi_fs = T.shared_roi_fs(dds)
    for i, dd in enumerate(dds, 1):
        T.render_deep_dive(prs, dd, brand, index=i, total=len(dds), roi_fs=roi_fs)

    if storyline.get("workforce_readiness"):
        T.render_delivery_capability(prs, storyline["workforce_readiness"], brand)

    T.render_approach(prs, brand)
    T.render_next_steps(prs, brand)

    idx = [-1] + [0] * (len(prs.slides) - 1)
    T.apply_uniform_footer(prs, meta.get("company", ""), ["deck"], idx, brand)
    prs.save(deck_path)
    say(f"Deck assembled — {len(prs.slides)} slides.")
    return {"deck_path": deck_path, "slide_count": len(prs.slides)}
