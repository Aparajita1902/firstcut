"""Phase 1 deck slide templates — rendered locally in Python (no LLM).

Each template adds one slide to a ``python-pptx`` ``Presentation`` object using
the supplied content and brand. All templates produce WHITE-background slides
with brand colours used as accents only (separator bars, section numbers,
italic subtitles). No dark theme anywhere.

Phase 1 covers four slide types:
    render_title_slide
    render_section_divider
    render_company_context
    render_next_steps

Plus a helper:
    apply_uniform_footer  — call ONCE at the end after all slides exist, to
                            write "Aistra × {Company} · {Section}" + "{n} / {N}"
                            on every slide in a consistent format.

Brand JSON lookups are defensive: every key has a fallback so the templates
still render if a brand key is missing.
"""

from __future__ import annotations

import re as _re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

import datetime as _datetime
import os as _os


# --------------------------------------------------------------------------------- #
# Brand defaults — used only when brand.json doesn't supply a key
# --------------------------------------------------------------------------------- #
_DEFAULTS = {
    "primary_purple": "6C48F2",
    "accent_cyan":    "00D9FF",
    "accent_pink":    "FF5CB6",
    "accent_teal":    "00E8FB",
    "ink":            "111827",
    "muted":          "64748B",
    "rule":           "E5E7EB",
    "card_fill":      "FAFAFA",
    "display_font":   "Georgia",
    "body_font":      "Calibri",
    "slide_width_inches":  13.333,
    "slide_height_inches":  7.5,
    "margin_inches":        0.5,
}


def _brand(brand: dict, key: str):
    """Read brand[key] if present, else fall back to _DEFAULTS[key]."""
    v = brand.get(key) if brand else None
    if v is None:
        return _DEFAULTS[key]
    return v.lstrip("#") if (isinstance(v, str) and key.endswith(("purple", "cyan", "pink", "teal", "ink", "muted", "rule", "fill"))) else v


def _rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str.lstrip("#"))


# --------------------------------------------------------------------------------- #
# Low-level slide-building helpers
# --------------------------------------------------------------------------------- #
def _add_blank_slide(prs: Presentation):
    """Add a blank slide (layout 6 is the standard 'Blank' layout) and force white bg."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Force background to solid white, even if the layout/master has fills set.
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return slide


def _add_textbox(slide, left, top, width, height, text, *,
                 font_name, font_size, bold=False, italic=False,
                 color_hex="111827", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Add a styled text box. Returns the textframe so callers can append runs."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text or ""
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color_hex)
    return tf


def _add_horizontal_line(slide, left, top, width, color_hex, thickness_pt=3):
    """Draw a horizontal accent line (a thin rectangle, so we control colour cleanly)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(thickness_pt))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color_hex)
    shape.line.fill.background()  # no border
    return shape


def _add_card(slide, left, top, width, height, fill_hex, border_hex=None, left_accent_hex=None):
    """Add a light card with optional 1pt border and/or 3pt left purple accent bar."""
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(fill_hex)
    if border_hex:
        card.line.color.rgb = _rgb(border_hex)
        card.line.width = Pt(1)
    else:
        card.line.fill.background()
    if left_accent_hex:
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(3), height)
        accent.fill.solid()
        accent.fill.fore_color.rgb = _rgb(left_accent_hex)
        accent.line.fill.background()
    return card


# --------------------------------------------------------------------------------- #
# Phase 1 templates
# --------------------------------------------------------------------------------- #
import math as _math


def _est_lines(text: str, font_size_pt: float, width_in: float, *, display: bool = False) -> int:
    """Estimate how many lines `text` wraps to in a box `width_in` wide at the given
    font size. Conservative (rounds up) so dynamic layouts leave enough room rather
    than overlap. `display` widens the per-char estimate slightly for serif headers."""
    if not text:
        return 1
    cw_factor = 0.50 if display else 0.48  # avg char advance as a fraction of em
    char_w_in = (font_size_pt / 72.0) * cw_factor
    cpl = max(1, int(width_in / char_w_in))
    return max(1, _math.ceil(len(str(text)) / cpl))


def _line_h_in(font_size_pt: float) -> float:
    """Single-line height in inches for a given point size (≈1.2 line spacing)."""
    return font_size_pt / 72.0 * 1.2


def _title_and_separator(slide, margin, content_w, title, *, display, purple, ink,
                         font_size=30, top_in=0.55):
    """Render a slide title and place the purple separator BELOW the actual title,
    so a title that wraps to two lines never collides with the bar. Returns the Y
    (in inches) at which following content should start."""
    lines = _est_lines(title, font_size, content_w / 914400.0, display=True)
    title_h = lines * _line_h_in(font_size)
    _add_textbox(slide, margin, Inches(top_in), content_w, Inches(title_h + 0.1),
                 title, font_name=display, font_size=font_size, color_hex=ink)
    sep_y = top_in + title_h + 0.13
    _add_horizontal_line(slide, margin, Inches(sep_y), Inches(0.6), purple, thickness_pt=3)
    return sep_y + 0.30  # where subtitle/body should begin
def _logo_path():
    """Resolve data/logo.png relative to this file (project root)."""
    p = _os.path.join(_os.path.dirname(_os.path.abspath(__file__)), "data", "logo.png")
    return p if _os.path.exists(p) else None


def _tint(hex_str: str, f: float = 0.86) -> str:
    h = hex_str.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = int(r + (255 - r) * f); g = int(g + (255 - g) * f); b = int(b + (255 - b) * f)
    return f"{r:02X}{g:02X}{b:02X}"


def _house_bars(slide, sw, sh, purple: str) -> None:
    for top in (Emu(0), sh - Pt(8)):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), top, sw, Pt(8))
        bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(purple)
        bar.line.fill.background(); bar.shadow.inherit = False


def _pill(slide, x_in, y_in, text, fill_hex, text_hex, body, *, size=13.33):
    w = min(4.4, max(1.7, len(text) * 0.095 + 0.55))
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_in), Inches(y_in), Inches(w), Inches(0.47))
    s.fill.solid(); s.fill.fore_color.rgb = _rgb(fill_hex)
    s.line.fill.background(); s.shadow.inherit = False
    try:
        s.adjustments[0] = 0.5
    except Exception:
        pass
    _add_textbox(slide, Inches(x_in), Inches(y_in - 0.005), Inches(w), Inches(0.47), text,
                 font_name=body, font_size=size, bold=True, color_hex=text_hex,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return w


_COLOR_TOKENS = {
    "ink": "111827", "purple": "6C48F2", "cyan": "0AB4D6", "pink": "EC4899",
    "green": "16A34A", "orange": "ED7D31", "slate": "475569", "muted": "94A3B8",
}


def _tok(color) -> str:
    """Resolve a semantic color token ('ink','purple',...) to hex; pass hex through."""
    if not isinstance(color, str):
        return "111827"
    return _COLOR_TOKENS.get(color.strip().lower(), color)


def _two_tone(slide, x_in, y_in, w_in, h_in, parts, body, *, size=37.33):
    tb = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = tb.text_frame; tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(0))
    p = tf.paragraphs[0]
    for text, color in parts:
        r = p.add_run(); r.text = text; r.font.name = body
        r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = _rgb(_tok(color))


def _bottom_logo(slide) -> None:
    lp = _logo_path()
    if lp:
        try:
            slide.shapes.add_picture(lp, Inches(0.53), Inches(6.82), height=Inches(0.42))
        except Exception:
            pass


def render_title_slide(prs: Presentation, title_content: dict, thesis: dict, brand: dict,
                       audience_role: str = "", date_str: str = "") -> None:
    """Title slide — 'AI Opportunities / for {company}', Aistra logo top-left, the purple
    circle motif filling the right, full-width purple bars top and bottom, and a month-year
    date. (audience_role / thesis are accepted for signature compatibility but not shown.)"""
    purple    = _brand(brand, "primary_purple")
    ink       = _brand(brand, "ink")
    muted     = _brand(brand, "muted")
    body      = _brand(brand, "body_font")
    outer_col = "7E61F3"   # lighter outer circle
    inner_col = purple     # primary (slightly darker) inner circle

    slide = _add_blank_slide(prs)
    sw = prs.slide_width
    sh = prs.slide_height
    sw_in = sw / 914400.0

    # --- RHS circle motif (drawn first; the text sits on the left) ------------- #
    outer_d = 5.6
    outer_left = sw_in - 4.0
    outer_top = -0.25
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(outer_left), Inches(outer_top),
                               Inches(outer_d), Inches(outer_d))
    o.fill.solid(); o.fill.fore_color.rgb = _rgb(outer_col); o.line.fill.background()
    o.shadow.inherit = False
    inner_d = 3.7
    inner_left = outer_left + (outer_d - inner_d) / 2 + 0.5
    inner_top = outer_top + (outer_d - inner_d) / 2 + 0.5
    ic = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(inner_left), Inches(inner_top),
                                Inches(inner_d), Inches(inner_d))
    ic.fill.solid(); ic.fill.fore_color.rgb = _rgb(inner_col); ic.line.fill.background()
    ic.shadow.inherit = False

    # --- Aistra logo, top-left ------------------------------------------------- #
    lp = _logo_path()
    placed = False
    if lp:
        try:
            slide.shapes.add_picture(lp, Inches(0.85), Inches(0.95), height=Inches(1.05))
            placed = True
        except Exception:
            placed = False
    if not placed:
        _add_textbox(slide, Inches(0.85), Inches(1.0), Inches(3), Inches(0.6),
                     "Aistra", font_name=body, font_size=28, bold=True, color_hex=purple)

    # --- Title: AI Opportunities / for {company} ------------------------------- #
    company = (title_content.get("title") or "Client Name").strip()
    _add_textbox(slide, Inches(0.85), Inches(2.95), Inches(8.2), Inches(0.95),
                 "AI Opportunities", font_name=body, font_size=42, bold=True, color_hex=ink)
    _add_textbox(slide, Inches(0.85), Inches(3.80), Inches(8.6), Inches(0.95),
                 f"for {company}", font_name=body, font_size=42, bold=True, color_hex=purple)

    # --- Descriptor: outside-in assessment ------------------------------------- #
    _add_textbox(slide, Inches(0.87), Inches(4.78), Inches(8.0), Inches(0.4),
                 "An outside-in assessment", font_name=body, font_size=20, color_hex=muted)

    # --- Date (Month, Year) ---------------------------------------------------- #
    if not (date_str or "").strip():
        date_str = _datetime.date.today().strftime("%B, %Y")
    _add_textbox(slide, Inches(0.85), Inches(5.45), Inches(6.0), Inches(0.4),
                 date_str, font_name=body, font_size=18, color_hex=muted)


def render_section_divider(prs: Presentation, number: str, name: str, transition: str,
                           brand: dict) -> None:
    """Section divider — large number in purple, name to the right, transitional italic.
    Right edge has a thin vertical purple accent bar for visual rhythm."""
    purple = _brand(brand, "primary_purple")
    ink    = _brand(brand, "ink")
    muted  = _brand(brand, "muted")
    display = _brand(brand, "display_font")
    body    = _brand(brand, "body_font")

    slide = _add_blank_slide(prs)
    sw = prs.slide_width
    sh = prs.slide_height
    margin = Inches(0.7)

    # Top label
    _add_textbox(slide, margin, Inches(0.5), Inches(4), Inches(0.3),
                 f"SECTION {number}", font_name=body, font_size=10, bold=True, color_hex=purple)

    # Large number — display font, purple
    _add_textbox(slide, margin, Inches(2.4), Inches(3.0), Inches(2.5),
                 number, font_name=display, font_size=140, color_hex=purple)

    # Section name to the right of the number
    _add_textbox(slide, Inches(4.0), Inches(2.7), sw - Inches(4.7), Inches(1.2),
                 name, font_name=display, font_size=40, color_hex=ink,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Italic transitional sentence below the name
    if transition:
        _add_textbox(slide, Inches(4.0), Inches(3.9), sw - Inches(4.7), Inches(2.0),
                     transition, font_name=body, font_size=14, italic=True, color_hex=muted)

    # Right-edge vertical accent bar (bookmark stripe) — thin purple line top to bottom
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        sw - Inches(0.25), Inches(0.7),
        Pt(3), sh - Inches(1.5),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = _rgb(purple)
    accent_bar.line.fill.background()


def _split_paragraphs(text: str, max_paras: int = 3) -> list[str]:
    """Split a wall-of-text body into 2–3 paragraphs for breathability."""
    text = (text or "").strip()
    if not text:
        return []
    # If the body already contains explicit paragraph breaks, respect them.
    if "\n\n" in text:
        paras = [p.strip() for p in text.split("\n\n") if p.strip()]
        return paras[:max_paras] if len(paras) <= max_paras else (
            paras[: max_paras - 1] + [" ".join(paras[max_paras - 1 :])]
        )
    # Otherwise split on sentence endings ('. ' followed by capital or start of next).
    import re as _re
    sentences = [s.strip() for s in _re.split(r"(?<=[.!?])\s+(?=[A-Z])", text) if s.strip()]
    if len(sentences) <= max_paras:
        return sentences
    # Group sentences evenly into ``max_paras`` paragraphs.
    per = max(1, len(sentences) // max_paras)
    groups = []
    for i in range(max_paras):
        start = i * per
        end = (i + 1) * per if i < max_paras - 1 else len(sentences)
        groups.append(" ".join(sentences[start:end]))
    return groups


def render_company_context(prs: Presentation, content: dict, brand: dict) -> None:
    """Company context slide. Title + 2–3 spaced paragraphs left; metrics card OR
    pull-quote card on the right. Designed to avoid the wall-of-text look."""
    purple = _brand(brand, "primary_purple")
    ink    = _brand(brand, "ink")
    muted  = _brand(brand, "muted")
    rule   = _brand(brand, "rule")
    card   = _brand(brand, "card_fill")
    display = _brand(brand, "display_font")
    body    = _brand(brand, "body_font")

    slide = _add_blank_slide(prs)
    sw = prs.slide_width
    margin = Inches(0.7)

    # Title + separator (separator sits below the actual title, so 2-line titles
    # don't collide with the bar). content_top is where body/panel begin.
    content_top_in = _title_and_separator(
        slide, margin, sw - 2 * margin, content.get("title") or "Company context",
        display=display, purple=purple, ink=ink, font_size=30, top_in=0.55,
    )
    content_top_in = max(content_top_in, 1.85)
    content_top = Inches(content_top_in)

    raw_body = (content.get("body") or "").strip()
    metrics = content.get("metrics") or []

    # If no metrics, derive a pull-quote from the body's most arresting sentence
    # (heuristic: the final sentence is typically the conclusion/punchline).
    paragraphs = _split_paragraphs(raw_body, max_paras=3)
    pull_quote = ""
    if not metrics and paragraphs:
        # Take the LAST sentence of the LAST paragraph as the pull-quote.
        last_para = paragraphs[-1]
        import re as _re
        sents = [s.strip() for s in _re.split(r"(?<=[.!?])\s+(?=[A-Z])", last_para) if s.strip()]
        if len(sents) >= 1 and sum(len(s) for s in sents) > 120:
            pull_quote = sents[-1]
            if len(sents) > 1:
                paragraphs[-1] = " ".join(sents[:-1])
            else:
                paragraphs = paragraphs[:-1]

    # Determine column widths
    has_right_panel = bool(metrics) or bool(pull_quote)
    if has_right_panel:
        body_w = Inches(7.3)
        right_left = margin + body_w + Inches(0.4)
        right_w = sw - right_left - margin
    else:
        body_w = sw - 2 * margin
        right_left = right_w = None

    # Render paragraphs as a single text frame with paragraph breaks — avoids
    # the dead-space-per-fixed-box artefact we'd get from multiple text boxes.
    if paragraphs:
        tb = slide.shapes.add_textbox(margin, content_top, body_w,
                                      Inches(max(2.0, 6.7 - content_top_in)))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(0)
        tf.margin_right = Pt(0)
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)
        tf.vertical_anchor = MSO_ANCHOR.TOP

        # First paragraph uses the default first paragraph
        first = tf.paragraphs[0]
        first.alignment = PP_ALIGN.LEFT
        first.space_after = Pt(10)
        r = first.add_run()
        r.text = paragraphs[0]
        r.font.name = body
        r.font.size = Pt(13)
        r.font.color.rgb = _rgb(ink)

        # Remaining paragraphs added with breathing-room spacing
        for para in paragraphs[1:]:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(10)
            p.space_after = Pt(10)
            r = p.add_run()
            r.text = para
            r.font.name = body
            r.font.size = Pt(13)
            r.font.color.rgb = _rgb(ink)

    # Right panel
    if metrics:
        # Existing metrics card behaviour
        right_top = content_top
        right_h = Inches(min(4.0, 0.7 + 0.7 * len(metrics)))
        _add_card(slide, right_left, right_top, right_w, right_h, card,
                  border_hex=rule, left_accent_hex=purple)
        _add_textbox(slide, right_left + Inches(0.25), right_top + Inches(0.2),
                     right_w - Inches(0.5), Inches(0.3),
                     (content.get("metrics_title") or "KEY FACTS").upper(),
                     font_name=body, font_size=10, bold=True, color_hex=purple)
        for i, m in enumerate(metrics):
            row_top = right_top + Inches(0.6 + 0.55 * i)
            value = str(m.get("value") or "")
            label = str(m.get("label") or "")
            _add_textbox(slide, right_left + Inches(0.25), row_top, Inches(1.4), Inches(0.5),
                         value, font_name=display, font_size=20, color_hex=purple)
            _add_textbox(slide, right_left + Inches(1.7), row_top + Inches(0.1),
                         right_w - Inches(1.95), Inches(0.4),
                         label, font_name=body, font_size=11, color_hex=ink)
    elif pull_quote:
        # Pull-quote card — large italic quote, prominent purple opening mark.
        right_top = content_top
        right_h = Inches(3.6)
        _add_card(slide, right_left, right_top, right_w, right_h, "FFFFFF",
                  border_hex=rule, left_accent_hex=purple)
        # Decorative opening quote mark
        _add_textbox(slide, right_left + Inches(0.25), right_top + Inches(0.05),
                     Inches(0.8), Inches(0.9), "\u201C",
                     font_name=display, font_size=54, color_hex=purple)
        # Quote text
        _add_textbox(slide, right_left + Inches(0.3), right_top + Inches(0.85),
                     right_w - Inches(0.6), right_h - Inches(1.0),
                     pull_quote, font_name=display, font_size=16, italic=True, color_hex=ink)


def render_next_steps(prs: Presentation, content: dict, brand: dict, company_name: str) -> None:
    """Next steps slide — two columns: WITH {COMPANY} | WITH AISTRA, numbered actions each."""
    purple = _brand(brand, "primary_purple")
    cyan   = _brand(brand, "accent_cyan")
    ink    = _brand(brand, "ink")
    muted  = _brand(brand, "muted")
    rule   = _brand(brand, "rule")
    card   = _brand(brand, "card_fill")
    display = _brand(brand, "display_font")
    body    = _brand(brand, "body_font")

    slide = _add_blank_slide(prs)
    sw = prs.slide_width
    margin = Inches(0.7)
    content_w = sw - 2 * margin

    # Title + separator (dynamic, clears 2-line titles)
    content_top_in = _title_and_separator(
        slide, margin, content_w, content.get("title") or "Next steps",
        display=display, purple=purple, ink=ink, font_size=30, top_in=0.55,
    )

    # Optional subtitle
    if content.get("subtitle"):
        _add_textbox(slide, margin, Inches(content_top_in), content_w, Inches(0.5),
                     content["subtitle"], font_name=body, font_size=13, italic=True, color_hex=muted)
        content_top_in += 0.55

    col_top_in = max(content_top_in + 0.2, 2.5)
    gap = Inches(0.4)
    col_w = (sw - 2 * margin - gap) / 2
    text_w_in = (col_w / 914400.0) - 1.15  # available width for action text, inches

    with_company = content.get("with_company") or content.get("client_actions") or []
    with_aistra  = content.get("with_aistra")  or content.get("aistra_actions") or []

    def _norm(action):
        return action if isinstance(action, str) else (action.get("action") or action.get("text") or "")

    def _layout(actions):
        """Return [(text, y_offset_in, row_h_in)] with variable row heights and the
        total stack height, so rows never overlap regardless of action length."""
        rows, y = [], 0.75  # start below the column header
        for action in actions[:5]:
            text = _norm(action)
            lines = _est_lines(text, 11, text_w_in)
            row_h = max(0.45, lines * _line_h_in(11) + 0.16)
            rows.append((text, y, row_h))
            y += row_h + 0.12
        return rows, y

    left_rows, left_total = _layout(with_company)
    right_rows, right_total = _layout(with_aistra)
    # One shared card height for symmetry; clamp so it never runs into the footer.
    col_h_in = min(6.85 - col_top_in, max(3.0, left_total, right_total) + 0.25)
    col_top = Inches(col_top_in)
    col_h = Inches(col_h_in)

    def _column(left, header_text, accent_hex, rows):
        _add_card(slide, left, col_top, col_w, col_h, card, border_hex=rule, left_accent_hex=accent_hex)
        _add_textbox(slide, left + Inches(0.3), col_top + Inches(0.2), col_w - Inches(0.6), Inches(0.4),
                     header_text.upper(), font_name=body, font_size=11, bold=True, color_hex=accent_hex)
        for i, (text, y_off, row_h) in enumerate(rows):
            row_top = col_top + Inches(y_off)
            _add_textbox(slide, left + Inches(0.3), row_top, Inches(0.5), Inches(0.4),
                         f"{i+1:02d}", font_name=display, font_size=18, color_hex=accent_hex)
            _add_textbox(slide, left + Inches(0.85), row_top + Inches(0.02),
                         col_w - Inches(1.15), Inches(row_h),
                         text, font_name=body, font_size=11, color_hex=ink)

    _column(margin, f"With {company_name}", purple, left_rows)
    _column(margin + col_w + gap, "With Aistra", cyan, right_rows)


def render_deep_dive(prs: Presentation, content: dict, brand: dict, *,
                     index: int | None = None, total: int | None = None,
                     roi_fs: float | None = None) -> None:
    """Deep-dive slide in the house style — DEEP-DIVE pill, area + position top-right,
    bold-sans opportunity title, THE CONTEXT (numbered) and THE OPPORTUNITY (icon levers)
    columns, and a purple INDICATIVE IMPACT ribbon. House chrome (bars, bottom logo, page
    number) is applied by apply_uniform_footer. Falls back to legacy what/why/how fields."""
    purple  = _brand(brand, "primary_purple")
    ink     = _brand(brand, "ink")
    muted   = _brand(brand, "muted")
    rule    = _brand(brand, "rule")
    body    = _brand(brand, "body_font")
    soft    = "EDE9FE"; white = "FFFFFF"; lav = "D9CEFF"; chipln = "B7A6FF"

    slide = _add_blank_slide(prs)
    sw = prs.slide_width
    sw_in = sw / 914400.0
    margin = Inches(0.7); margin_in = 0.7
    content_w_in = sw_in - 1.4

    # Pill + area/position meta
    _pill(slide, 0.7, 0.5, "Deep-dive", purple, white, body, size=13.33)
    area = (content.get("area") or "").strip()
    meta = []
    if area:
        meta.append(area.upper())
    if index and total:
        meta.append(f"{index} OF {total}")
    if meta:
        _add_textbox(slide, Inches(7.0), Inches(0.62), Inches(sw_in - margin_in - 7.0), Inches(0.3),
                     "   \u00b7   ".join(meta), font_name=body, font_size=10.5, bold=True,
                     color_hex=muted, align=PP_ALIGN.RIGHT)

    # Title (bold sans)
    title = (content.get("title") or "Opportunity").strip()
    t_fs = 30
    tlines = _est_lines(title, t_fs, content_w_in)
    t_top = 1.12
    _add_textbox(slide, margin, Inches(t_top), Inches(content_w_in), Inches(tlines * _line_h_in(t_fs) + 0.1),
                 title, font_name=body, font_size=t_fs, bold=True, color_hex=ink)
    y = t_top + tlines * _line_h_in(t_fs) + 0.04

    # Subtitle
    subtitle = (content.get("subtitle") or "").strip()
    if subtitle:
        s_fs = 13
        slines = _est_lines(subtitle, s_fs, content_w_in)
        _add_textbox(slide, margin, Inches(y), Inches(content_w_in), Inches(slines * _line_h_in(s_fs) + 0.1),
                     subtitle, font_name=body, font_size=s_fs, italic=True, color_hex=muted)
        y += slines * _line_h_in(s_fs) + 0.12
    body_top = y + 0.18

    # Geometry: two columns + bottom ribbon (ribbon nudged up to clear the bottom logo)
    ribbon_top = 5.45; ribbon_h = 1.05; col_bottom = ribbon_top - 0.42
    left_w_in = 5.35; gap_in = 0.55
    right_x_in = margin_in + left_w_in + gap_in
    right_w_in = sw_in - margin_in - right_x_in
    right_x = Inches(right_x_in)

    vdiv_x = Inches(margin_in + left_w_in + gap_in / 2)
    vdiv = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, vdiv_x, Inches(body_top), Pt(1), Inches(col_bottom - body_top))
    vdiv.fill.solid(); vdiv.fill.fore_color.rgb = _rgb(rule); vdiv.line.fill.background(); vdiv.shadow.inherit = False

    # LEFT: THE CONTEXT
    _add_textbox(slide, margin, Inches(body_top), Inches(left_w_in), Inches(0.26),
                 "THE CONTEXT", font_name=body, font_size=12, bold=True, color_hex=muted)
    pts = content.get("context_points")
    if not pts:
        pts = [content.get("what"), content.get("why_this_company"), content.get("so_what")]
    pts = [str(p).strip() for p in (pts or []) if str(p).strip()][:5]

    levers = content.get("opportunity_levers")
    if not levers:
        legacy = content.get("how_aistra_delivers")
        levers = [legacy] if (legacy or "").strip() else []
    levers = [str(l).strip() for l in (levers or []) if str(l).strip()][:4]

    # Shrink-with-floor so both columns clear the ribbon, whichever is denser.
    txt_w_in = left_w_in - 0.42
    lev_txt_w_in = right_w_in - 0.52
    col_avail = col_bottom - (body_top + 0.42)

    def _col_h(items, width_in, fs, min_row, gap):
        h = 0.0
        for it in items:
            h += max(min_row, _est_lines(it, fs, width_in) * _line_h_in(fs)) + gap
        return h

    fs_fit = 14
    while fs_fit > 11 and (
        _col_h(pts, txt_w_in, fs_fit, 0.34, 0.16) > col_avail
        or _col_h(levers, lev_txt_w_in, fs_fit, 0.40, 0.18) > col_avail
    ):
        fs_fit -= 0.5
    p_fs = l_fs = fs_fit

    txt_x = margin + Inches(0.42)
    cy = body_top + 0.42
    for n, pt in enumerate(pts, start=1):
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, margin, Inches(cy), Inches(0.26), Inches(0.26))
        circ.fill.background(); circ.line.color.rgb = _rgb(rule); circ.line.width = Pt(1.25); circ.shadow.inherit = False
        _add_textbox(slide, margin, Inches(cy), Inches(0.26), Inches(0.26), str(n),
                     font_name=body, font_size=9, bold=True, color_hex=muted,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        nlines = _est_lines(pt, p_fs, txt_w_in)
        rh = nlines * _line_h_in(p_fs)
        _add_textbox(slide, txt_x, Inches(cy - 0.02), Inches(txt_w_in), Inches(rh + 0.1),
                     pt, font_name=body, font_size=p_fs, color_hex=ink)
        cy += max(0.34, rh) + 0.16

    # RIGHT: THE OPPORTUNITY
    _add_textbox(slide, right_x, Inches(body_top), Inches(right_w_in), Inches(0.26),
                 "THE OPPORTUNITY", font_name=body, font_size=12, bold=True, color_hex=purple)
    lev_txt_x = right_x + Inches(0.52)
    ry = body_top + 0.42
    for lev in levers:
        tile = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, Inches(ry), Inches(0.36), Inches(0.36))
        tile.fill.solid(); tile.fill.fore_color.rgb = _rgb(soft); tile.line.fill.background(); tile.shadow.inherit = False
        try:
            tile.adjustments[0] = 0.28
        except Exception:
            pass
        dia = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, right_x + Inches(0.12), Inches(ry + 0.12), Inches(0.12), Inches(0.12))
        dia.fill.solid(); dia.fill.fore_color.rgb = _rgb(purple); dia.line.fill.background(); dia.shadow.inherit = False
        nlines = _est_lines(lev, l_fs, lev_txt_w_in)
        rh = nlines * _line_h_in(l_fs)
        _add_textbox(slide, lev_txt_x, Inches(ry - 0.01), Inches(lev_txt_w_in), Inches(rh + 0.1),
                     lev, font_name=body, font_size=l_fs, color_hex=ink)
        ry += max(0.40, rh) + 0.18

    # BOTTOM RIBBON
    band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, margin, Inches(ribbon_top),
                                  Inches(content_w_in), Inches(ribbon_h))
    band.fill.solid(); band.fill.fore_color.rgb = _rgb(purple); band.line.fill.background(); band.shadow.inherit = False
    try:
        band.adjustments[0] = 0.10
    except Exception:
        pass
    roi = content.get("roi") or {}
    low = str(roi.get("low") or "").strip(); high = str(roi.get("high") or "").strip()
    pad = 0.34; inL = margin_in + pad
    right_edge_in = sw_in - margin_in - pad

    # Number lives in a fixed left zone; chips in a bounded right zone with a hard gap
    # between them so they can never collide regardless of length.
    num_zone_w = 3.5
    zone_gap = 0.5
    chips_zone_left = inL + num_zone_w + zone_gap
    chips_zone_w = right_edge_in - chips_zone_left

    _add_textbox(slide, Inches(inL), Inches(ribbon_top + 0.22), Inches(num_zone_w), Inches(0.22),
                 "INDICATIVE IMPACT \u00b7 ANNUAL", font_name=body, font_size=9, bold=True, color_hex=lav)
    if _has_digit(low) or _has_digit(high):
        rng = _collapse_range(low, high)
        num_fs = roi_fs if roi_fs else _fit_one_line_fs(rng, num_zone_w, 26, 16)
        _add_textbox(slide, Inches(inL), Inches(ribbon_top + 0.44), Inches(num_zone_w), Inches(0.55),
                     rng, font_name=body, font_size=num_fs, bold=True, color_hex=white)
    else:
        # Non-numeric ROI -> small label, never the giant-number treatment.
        _add_textbox(slide, Inches(inL), Inches(ribbon_top + 0.50), Inches(num_zone_w), Inches(0.4),
                     "Directional", font_name=body, font_size=18, bold=True, color_hex=white)

    drivers = [d for d in (roi.get("drivers") or []) if (d.get("name") or d.get("value"))][:3]
    _add_textbox(slide, Inches(right_edge_in - chips_zone_w), Inches(ribbon_top + 0.22),
                 Inches(chips_zone_w), Inches(0.22),
                 "IMPACT LEVERS", font_name=body, font_size=9, bold=True, color_hex=lav,
                 align=PP_ALIGN.RIGHT)
    # Name the levers only — short labels, capped in the prompt so they fit without truncation.
    chip_texts = [(d.get("name") or d.get("value") or "").strip() for d in drivers]
    chip_texts = [t for t in chip_texts if t]

    # Shrink chip font (to a floor) until the row fits the right zone; then clamp text.
    chip_fs = 11
    def _chip_widths(fs):
        cw_ = fs / 72.0 * 0.52
        return [max(1.0, len(t) * cw_ + 0.34) for t in chip_texts]
    gap = 0.14
    widths = _chip_widths(chip_fs)
    while chip_fs > 8.5 and (sum(widths) + gap * (len(widths) - 1 if widths else 0)) > chips_zone_w:
        chip_fs -= 0.5
        widths = _chip_widths(chip_fs)
    # If still too wide at the floor, clamp each chip's text to its fair share of the zone.
    if chip_texts and (sum(widths) + gap * (len(widths) - 1)) > chips_zone_w:
        per = (chips_zone_w - gap * (len(chip_texts) - 1)) / len(chip_texts)
        chip_texts = [_clamp(t, per - 0.34, chip_fs, 1) for t in chip_texts]
        widths = _chip_widths(chip_fs)
    total_w = sum(widths) + gap * (len(widths) - 1 if widths else 0)
    cx = max(chips_zone_left, right_edge_in - total_w)
    chip_y = ribbon_top + 0.50; chip_h = 0.36
    for t, w in zip(chip_texts, widths):
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(chip_y), Inches(w), Inches(chip_h))
        pill.fill.background(); pill.line.color.rgb = _rgb(chipln); pill.line.width = Pt(1); pill.shadow.inherit = False
        try:
            pill.adjustments[0] = 0.5
        except Exception:
            pass
        _add_textbox(slide, Inches(cx), Inches(chip_y), Inches(w), Inches(chip_h), t,
                     font_name=body, font_size=chip_fs, color_hex=white,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w + gap


# --------------------------------------------------------------------------------- #
# Footer + uniform page numbers (apply once at end)
# --------------------------------------------------------------------------------- #
def apply_uniform_footer(prs: Presentation, company_name: str, section_names: list[str],
                          slide_section_index: list[int], brand: dict) -> None:
    """Apply house chrome to every slide: top/bottom purple bars, the Aistra logo bottom-left
    (content slides only; the title carries its own top-left logo), and a plain page number
    bottom-right. ``slide_section_index`` marks the title slide with -1. (``company_name`` and
    ``section_names`` are kept for signature compatibility.)"""
    purple = _brand(brand, "primary_purple")
    muted  = _brand(brand, "muted")
    body   = _brand(brand, "body_font")
    sw = prs.slide_width
    sh = prs.slide_height
    for i, slide in enumerate(prs.slides, start=1):
        idx = slide_section_index[i - 1] if (i - 1) < len(slide_section_index) else -1
        _house_bars(slide, sw, sh, purple)
        if idx >= 0:
            _bottom_logo(slide)
        _add_textbox(slide, sw - Inches(1.05), sh - Inches(0.62), Inches(0.6), Inches(0.3),
                     str(i), font_name=body, font_size=12, color_hex=muted, align=PP_ALIGN.RIGHT)


# =================================================================================== #
# House-style content renderers (chrome — bars/logo/page number — added by
# apply_uniform_footer). Ported from the aligned standalone mocks.
# =================================================================================== #
_PILLARS = {
    "revenue":     "6C48F2",
    "customer":    "0AB4D6",
    "operational": "ED7D31",
    "financial":   "EC4899",
    "workforce":   "16A34A",
}
_PILLAR_LABEL = {
    "revenue": "Revenue", "customer": "Customer", "operational": "Operational",
    "financial": "Financial", "workforce": "Workforce",
}
# Representative icon per pillar (for opportunity-areas cards)
_PILLAR_ICON = {
    "revenue": "revenue_crosssell", "customer": "customer_conversational",
    "operational": "operational_docs", "financial": "financial_collections",
    "workforce": "workforce_planning",
}


def _t(slide, x, y, w, h, text, **kw):
    return _add_textbox(slide, Inches(x), Inches(y), Inches(w), Inches(h), text, **kw)


# --------------------------------------------------------------------------------- #
# Rendering robustness helpers (cosmetics pass)
# --------------------------------------------------------------------------------- #
_CANON_PILLARS = ("revenue", "customer", "operational", "financial", "workforce")


def _norm_pillar(p) -> str:
    """Normalise any pillar string the model emits to a canonical lowercase id.

    Tolerates title case, trailing 'Intelligence', and stray whitespace, e.g.
    'Revenue', 'Revenue Intelligence', ' OPERATIONAL ' -> 'revenue'/'operational'.
    Falls back to 'customer' if nothing matches.
    """
    if not isinstance(p, str):
        return "customer"
    s = p.strip().lower().replace("intelligence", "").strip()
    if s in _CANON_PILLARS:
        return s
    for c in _CANON_PILLARS:
        if c in s:
            return c
    return "customer"


def _collapse_range(low: str, high: str) -> str:
    """'INR 60 Cr' + 'INR 150 Cr' -> 'INR 60-150 Cr'; '$1.5m' + '$4.0m' -> '$1.5-4.0m'.

    Factors a shared non-numeric prefix (currency) and suffix (unit) when both ends
    share them; otherwise falls back to 'low - high'. Returns low or high alone if
    only one is present.
    """
    low = (low or "").strip(); high = (high or "").strip()
    if not (low and high):
        return low or high or "\u2014"
    m1 = _re.match(r"^(\D*)([\d.,]+)(\D*)$", low)
    m2 = _re.match(r"^(\D*)([\d.,]+)(\D*)$", high)
    if m1 and m2 and m1.group(1) == m2.group(1) and m1.group(3) == m2.group(3):
        return f"{m1.group(1)}{m1.group(2)}\u2013{m2.group(2)}{m1.group(3)}"
    return f"{low} \u2013 {high}"


def _has_digit(s: str) -> bool:
    return any(ch.isdigit() for ch in (s or ""))


def shared_roi_fs(deep_dives: list, num_zone_w: float = 3.5, start_fs: float = 26,
                  floor_fs: float = 16) -> float:
    """One ROI number font size for every deep-dive, so the figure is the same size on
    every page. Sized to the LONGEST collapsed range across the deck so none overflows."""
    longest = ""
    for dd in deep_dives or []:
        roi = (dd or {}).get("roi") or {}
        low = str(roi.get("low") or "").strip(); high = str(roi.get("high") or "").strip()
        if _has_digit(low) or _has_digit(high):
            rng = _collapse_range(low, high)
            if len(rng) > len(longest):
                longest = rng
    return _fit_one_line_fs(longest, num_zone_w, start_fs, floor_fs) if longest else start_fs


def _fit_one_line_fs(text: str, width_in: float, start_fs: float, floor_fs: float,
                     *, bold: bool = True) -> float:
    """Largest font size (start..floor) at which `text` fits one line in width_in."""
    if not text:
        return start_fs
    fs = start_fs
    while fs > floor_fs and _est_lines(text, fs, width_in) > 1:
        fs -= 0.5
    return fs


def _clamp(text: str, width_in: float, fs: float, max_lines: int) -> str:
    """Char-clamp `text` so it renders in <= max_lines at fs/width; ellipsis if cut."""
    text = (text or "").strip()
    if not text or _est_lines(text, fs, width_in) <= max_lines:
        return text
    per_line = max(8, len(text) // max(1, _est_lines(text, fs, width_in)))
    budget = max(8, per_line * max_lines - 1)
    return text[:budget].rstrip(" ,;:") + "\u2026"


def _rrect(slide, x, y, w, h, fill, *, adj=0.5, line=None, line_w=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = _rgb(fill)
    if line:
        s.line.color.rgb = _rgb(line); s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    try:
        s.adjustments[0] = adj
    except Exception:
        pass
    s.shadow.inherit = False
    return s


def _rect2(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = _rgb(fill)
    s.line.fill.background(); s.shadow.inherit = False
    return s


def _icons_dir():
    return _os.path.join(_os.path.dirname(_os.path.abspath(__file__)), "data", "icons")


def _icon_tile(slide, x, y, d, pillar_hex, icon_name, *, muted=False):
    """Tinted circle + glyph PNG centered. Uses *_muted.png when muted."""
    circ_hex = "E5E7EB" if muted else _tint(pillar_hex, 0.86)
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = _rgb(circ_hex)
    c.line.fill.background(); c.shadow.inherit = False
    fn = icon_name + ("_muted" if muted else "") + ".png"
    p = _os.path.join(_icons_dir(), fn)
    if _os.path.exists(p):
        gs = d * 0.52
        try:
            slide.shapes.add_picture(p, Inches(x + (d - gs) / 2), Inches(y + (d - gs) / 2),
                                     Inches(gs), Inches(gs))
        except Exception:
            pass


def _header(slide, pill_text, title_parts, brand, *, subtitle="",
            subtitle_color=None, subtitle_italic=False):
    """Pill + two-tone title + optional subtitle. Returns the title height (0.93 or 1.48)."""
    purple = _brand(brand, "primary_purple")
    body = _brand(brand, "body_font")
    _pill(slide, 0.8, 0.53, pill_text, purple, "FFFFFF", body)
    n = sum(len(t) for t, _ in title_parts)
    th = 1.48 if n > 52 else 0.93
    _two_tone(slide, 0.8, 1.27, 11.9, th, title_parts, body, size=37.33)
    if subtitle:
        col = subtitle_color or _brand(brand, "muted")
        _t(slide, 0.82, 1.27 + th + 0.04, 11.6, 0.45, subtitle,
           font_name=body, font_size=15.5, italic=subtitle_italic, color_hex=col)
    return th


def _sources(slide, text, brand):
    if text:
        _t(slide, 1.6, 6.52, 10.2, 0.3, text, font_name="Arial", font_size=10.67,
           italic=True, color_hex=_brand(brand, "muted"))


# --------------------------------------------------------------------------------- #
# Strategic context — 3/4 metric cards, or two-column compare
# --------------------------------------------------------------------------------- #
_SC_ACCENTS = [  # bar, metric, tag-text
    ("EC4899", "EC4899", "EC4899"),
    ("00D9FF", "00D9FF", "0891B2"),
    ("D97706", "ED7D31", "D97706"),
    ("00B050", "00B050", "00B050"),
]


def render_strategic_context_cards(prs, content, brand):
    ink = _brand(brand, "ink"); slate = _brand(brand, "muted")
    slide = _add_blank_slide(prs)
    two_line = sum(len(t) for t, _ in content.get("title_parts", [])) > 46
    _header(slide, content.get("pill", "Strategic context"), content.get("title_parts", [(" ", "ink")]), brand,
            subtitle=("" if two_line else content.get("subtitle", "")),
            subtitle_color=_brand(brand, "muted"), subtitle_italic=True)
    cards = content.get("cards", [])
    n = len(cards)
    if n <= 3:
        xs = [0.8, 4.93, 9.07][:n]; cw = 3.73; msz = 48; cy_m = 3.33; cy_d = 4.30; cy_b = 4.95
    else:
        xs = [0.8, 3.84, 6.88, 10.07][:n]; cw = 2.85; msz = 37.33; cy_m = 3.14; cy_d = 3.96; cy_b = 4.55
    two_line2 = sum(len(t) for t, _ in content.get("title_parts", [])) > 46
    dy = 0.30 if two_line2 else 0.0
    cy = 2.28 + dy; ch = 4.12 - dy
    body_bottom = cy + ch - 0.18
    for x, c, acc in zip(xs, cards, _SC_ACCENTS):
        bar, met, tag = acc
        _rrect(slide, x, cy, cw, ch, "F1F0FB", adj=0.06)
        _rect2(slide, x, cy, cw, 0.07, bar)
        tagw = min(cw - 0.4, max(1.4, len(c.get("tag", "")) * 0.085 + 0.4))
        _rrect(slide, x + 0.33, 2.55 + dy, tagw, 0.37, _tint(tag), adj=0.5)
        _t(slide, x + 0.33, 2.545 + dy, tagw, 0.37, c.get("tag", ""), font_name="Calibri", font_size=10.67,
           bold=True, color_hex=tag, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _t(slide, x + 0.33, cy_m + dy, cw - 0.5, 0.8, c.get("metric", ""), font_name="Calibri",
           font_size=msz, bold=True, color_hex=met)
        _t(slide, x + 0.33, cy_d + dy, cw - 0.6, 0.55, c.get("desc", ""), font_name="Calibri",
           font_size=14.67, bold=True, color_hex=ink)
        body_top = cy_b + dy
        body_h = max(0.6, body_bottom - body_top)
        body_lines = max(2, int(body_h / _line_h_in(14)))
        _t(slide, x + 0.33, body_top, cw - 0.6, body_h,
           _clamp(c.get("body", ""), cw - 0.6, 14, body_lines),
           font_name="Calibri Light", font_size=14, color_hex=slate)
    _sources(slide, content.get("sources", ""), brand)


def render_strategic_context_compare(prs, content, brand):
    slate = _brand(brand, "muted")
    slide = _add_blank_slide(prs)
    subtitle = content.get("subtitle", "")
    th = _header(slide, content.get("pill", "Strategic context"),
                 content.get("title_parts", [(" ", "ink")]), brand,
                 subtitle=subtitle, subtitle_color=_brand(brand, "muted"),
                 subtitle_italic=True)

    # Start the cards BELOW the (possibly two-line) subtitle so nothing bleeds through.
    sub_h = (_est_lines(subtitle, 15.5, 11.6) * _line_h_in(15.5) + 0.10) if subtitle else 0.0
    cy = max(2.55, 1.27 + th + 0.06 + sub_h + 0.12)
    bottom = 6.40
    ch = bottom - cy

    cols = [(0.8, "6C48F2", "16A34A", "\u2713", content.get("left", {"header": "", "items": []})),
            (7.07, "EC4899", "EC4899", "\u2717", content.get("right", {"header": "", "items": []}))]
    cw = 5.47
    item_w = cw - 1.1
    items_top = cy + 0.66
    avail = (cy + ch - 0.14) - items_top
    gap = 0.10
    MAXL = 2  # compare bullets clamp to two lines so columns stay scannable and bounded

    # One shared font size that fits the denser of the two columns (keeps them matched).
    def _col_items(col):
        return [_clamp(str(it).strip(), item_w, 13, MAXL)
                for it in (col.get("items") or [])[:5] if str(it).strip()]
    def _col_h(items, fs):
        return sum(_est_lines(it, fs, item_w) * _line_h_in(fs) + gap for it in items)
    fs = 13.0
    while fs > 9.0 and max(_col_h(_col_items(c[4]), fs) for c in cols) > avail:
        fs -= 0.5

    for x, hd, ic_col, ic, col in cols:
        _rrect(slide, x, cy, cw, ch, "F1F0FB", adj=0.05)
        _rect2(slide, x, cy, cw, 0.07, hd)
        _t(slide, x + 0.33, cy + 0.18, cw - 0.6, 0.4, col.get("header", ""), font_name="Calibri",
           font_size=16, bold=True, color_hex=hd)
        ry = items_top
        for it in _col_items(col):
            rh = _est_lines(it, fs, item_w) * _line_h_in(fs)
            _t(slide, x + 0.33, ry, 0.33, 0.30, ic, font_name="Calibri", font_size=min(15, fs + 1.5),
               bold=True, color_hex=ic_col, align=PP_ALIGN.CENTER)
            _t(slide, x + 0.73, ry, item_w, rh + 0.12, it, font_name="Calibri Light",
               font_size=fs, color_hex=slate)
            ry += rh + gap
    _sources(slide, content.get("sources", ""), brand)


# --------------------------------------------------------------------------------- #
# Art of the possible — fixed 5-pillar radial (10 buckets)
# --------------------------------------------------------------------------------- #
_AOTP = [
    # id, pillar, icon, label, x, y, edge
    ("revenue_pricing",            "revenue",     "revenue_pricing",            "Dynamic pricing & yield",               3.875, 2.656, "eb"),
    ("revenue_crosssell",          "revenue",     "revenue_crosssell",          "Next-best-action cross-sell",           6.75,  2.656, "eb"),
    ("customer_conversational",    "customer",    "customer_conversational",    "Conversational service & journeys",     0.604, 3.521, "er"),
    ("customer_churn",             "customer",    "customer_churn",             "Churn & voice-of-customer signals",     0.604, 4.4375,"er"),
    ("workforce_planning",         "workforce",   "workforce_planning",         "Demand-matched workforce planning",     0.604, 5.354, "er"),
    ("operational_docs",           "operational", "operational_docs",           "Document & process automation",         10.125,3.521, "el"),
    ("operational_orchestration",  "operational", "operational_orchestration",  "Intelligent orchestration & recovery",  10.125,4.4375,"el"),
    ("financial_collections",      "financial",   "financial_collections",      "Collections optimisation",              10.125,5.354, "el"),
    ("workforce_knowledge",        "workforce",   "workforce_knowledge",        "Frontline knowledge assistant",         3.875, 6.167, "et"),
    ("financial_fraud",            "financial",   "financial_fraud",            "Anomaly & fraud detection",             6.75,  6.167, "et"),
]


def render_art_of_possible(prs, content, brand):
    ink = _brand(brand, "ink"); purple = _brand(brand, "primary_purple")
    slate = _brand(brand, "muted"); body = _brand(brand, "body_font")
    company = content.get("company", "the business")
    flags = content.get("flags", {})
    slide = _add_blank_slide(prs)
    parts = [("What a sharper ", ink), (company, purple), (" could look like.", ink)]
    tot = sum(len(t) for t, _ in parts)
    tsize = 37.33 if tot <= 46 else 30
    _pill(slide, 0.8, 0.53, "Art of the possible", purple, "FFFFFF", body)
    _two_tone(slide, 0.8, 1.18, 11.9, 1.0, parts, body, size=tsize)

    # Legend (5 pillars), centered
    order = ["revenue", "customer", "operational", "financial", "workforce"]
    seg = []
    for p in order:
        seg.append((_PILLAR_LABEL[p], _PILLARS[p]))
    total_w = sum(0.16 + len(lbl) * 0.085 + 0.34 for lbl, _ in seg)
    lx = (13.333 - total_w) / 2
    ly = 2.12
    for lbl, col in seg:
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(lx), Inches(ly + 0.04), Inches(0.12), Inches(0.12))
        d.fill.solid(); d.fill.fore_color.rgb = _rgb(col); d.line.fill.background(); d.shadow.inherit = False
        _t(slide, lx + 0.18, ly - 0.02, len(lbl) * 0.085 + 0.2, 0.26, lbl,
           font_name=body, font_size=10.5, bold=True, color_hex=slate)
        lx += 0.16 + len(lbl) * 0.085 + 0.34

    # Center node
    node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.0), Inches(3.76), Inches(1.333), Inches(1.333))
    node.fill.solid(); node.fill.fore_color.rgb = _rgb(ink); node.line.fill.background(); node.shadow.inherit = False
    _t(slide, 6.0, 4.05, 1.333, 0.75, company.upper(), font_name=_brand(brand, "display_font"),
       font_size=13, color_hex="FFFFFF", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    cw, chh = 2.604, 0.667
    positions = [(x, y, edge) for (_id, _p, _ic, _lbl, x, y, edge) in _AOTP]
    cards = content.get("cards")
    items = []
    if cards:
        for i, card in enumerate(cards[:10]):
            x, y, edge = positions[i]
            pillar = card.get("pillar", "customer")
            items.append({"pillar": pillar,
                          "icon": card.get("icon") or _PILLAR_ICON.get(pillar, "customer_conversational"),
                          "label": card.get("label", ""),
                          "relevance": card.get("relevance", "core"),
                          "deep_dive": bool(card.get("deep_dive")),
                          "x": x, "y": y, "edge": edge})
    else:
        for (_id, pillar, icon, label, x, y, edge) in _AOTP:
            fl = flags.get(_id, {})
            items.append({"pillar": pillar, "icon": icon, "label": label,
                          "relevance": fl.get("relevance", "core"),
                          "deep_dive": bool(fl.get("deep_dive")),
                          "x": x, "y": y, "edge": edge})

    for it in items:
        pillar = it["pillar"]; x = it["x"]; y = it["y"]; edge = it["edge"]
        out = (it["relevance"] == "out")
        pcol = _PILLARS.get(pillar, "6C48F2")
        col = "9CA3AF" if out else pcol
        _rrect(slide, x, y, cw, chh, "FFFFFF", adj=0.20, line=col, line_w=1.6)
        if edge == "er":
            _rect2(slide, x + cw - 0.05, y + 0.12, 0.05, chh - 0.24, col)
        elif edge == "el":
            _rect2(slide, x, y + 0.12, 0.05, chh - 0.24, col)
        elif edge == "et":
            _rect2(slide, x + 0.18, y, cw - 0.36, 0.05, col)
        elif edge == "eb":
            _rect2(slide, x + 0.18, y + chh - 0.05, cw - 0.36, 0.05, col)
        _icon_tile(slide, x + 0.14, y + (chh - 0.42) / 2, 0.42, pcol, it["icon"], muted=out)
        _t(slide, x + 0.7, y + 0.04, cw - 0.84, chh - 0.08, it["label"], font_name=body,
           font_size=10.5, bold=True, color_hex=("9CA3AF" if out else ink), anchor=MSO_ANCHOR.MIDDLE)
        if it["deep_dive"]:
            m = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + cw - 0.2), Inches(y - 0.07), Inches(0.16), Inches(0.16))
            m.fill.solid(); m.fill.fore_color.rgb = _rgb(pcol); m.line.color.rgb = _rgb("FFFFFF")
            m.line.width = Pt(1.5); m.shadow.inherit = False


# --------------------------------------------------------------------------------- #
# Opportunity areas — 4 (2x2) or 6 (2x3) cards
# --------------------------------------------------------------------------------- #
def render_opportunity_areas(prs, content, brand):
    ink = _brand(brand, "ink"); slate = _brand(brand, "muted"); body = _brand(brand, "body_font")
    opps = content.get("opportunities", [])
    n = len(opps)
    slide = _add_blank_slide(prs)
    th = _header(slide, "Opportunity areas", content.get("title_parts", [(" ", "ink")]), brand,
                 subtitle="", subtitle_color=_brand(brand, "muted"),
                 subtitle_italic=True)
    top = 1.27 + th + 0.42
    bottom = 6.55
    xs = [0.667, 6.917]; cw = 5.75
    if n <= 4:
        gap = 0.18; ch = (bottom - top - gap) / 2
        ys = [top, top + ch + gap]; circ, tisz, bosz = 0.6, 15, 11
    else:
        gap = 0.16; ch = (bottom - top - 2 * gap) / 3
        ys = [top, top + ch + gap, top + 2 * (ch + gap)]; circ, tisz, bosz = 0.5, 13, 10
    pos = [(x, y) for y in ys for x in xs][:n]
    for (x, y), o in zip(pos, opps):
        pillar = _norm_pillar(o.get("pillar", "customer"))
        col = _PILLARS.get(pillar, "6C48F2")
        icon = o.get("icon") or _PILLAR_ICON.get(pillar, "customer_conversational")
        _rrect(slide, x, y, cw, ch, "FAFAFC", adj=0.07, line="ECECF3", line_w=1.0)
        _rect2(slide, x, y, 0.055, ch, col)  # left accent
        iy = y + (ch - circ) / 2
        _icon_tile(slide, x + 0.28, iy, circ, col, icon)
        tx = x + 0.28 + circ + 0.22
        tw = cw - (0.28 + circ + 0.22) - 0.2
        ty = y + 0.22
        _t(slide, tx, ty, tw, 0.2, _PILLAR_LABEL.get(pillar, "").upper() + " INTELLIGENCE",
           font_name=body, font_size=8, bold=True, color_hex=col)
        _t(slide, tx, ty + 0.2, tw, 0.4, o.get("title", ""), font_name=body, font_size=tisz,
           bold=True, color_hex=ink)
        body_top = ty + 0.2 + 0.36
        body_h = max(0.4, (y + ch) - body_top - 0.10)
        body_lines = max(1, int(body_h / _line_h_in(bosz)))
        _t(slide, tx, body_top, tw, body_h,
           _clamp(o.get("body", ""), tw * 0.94, bosz, body_lines),
           font_name="Calibri Light", font_size=bosz, color_hex=slate)


# --------------------------------------------------------------------------------- #
# Our Approach (static) + Next Steps (static) — generic, client-free copy
# --------------------------------------------------------------------------------- #
_APPROACH_STAGES = [
    ("01", "Week 1", "Discovery",
     "Focused working sessions with your leadership to map the current state, agree on the priority use case for the prototype, and confirm data requirements."),
    ("02", "Week 2", "Scope & Terms",
     "A short, precise scope document \u2014 selected use case, data and integration requirements, timeline, and commercial terms. No lengthy proposals."),
    ("03", "Weeks 3\u20136", "Prototype",
     "A working prototype built on real, anonymised data and validated against your business rules, then demonstrated to your leadership team."),
    ("04", "Weeks 7\u201315", "Production",
     "The approved prototype goes live within a limited, controlled scope, then rolls out to further areas on a structured programme."),
    ("05", "Ongoing", "Operate & Improve",
     "Aistra-managed operations with SLA-backed delivery, continuous tuning, quarterly reviews, and a roadmap for incremental opportunity areas."),
]
_NEXT_STEPS = [
    ("01", "Confirm priority focus",
     "Alignment on the solution areas \u2014 and which one we lead the prototype with.",
     "Your view on priority sequencing."),
    ("02", "Identify the working group",
     "Names from your side for the focused working session; Aistra brings the engagement team.",
     "Attendees and a 90-minute slot within 10 days."),
    ("03", "Agree data access for the prototype",
     "Anonymised data relevant to the chosen use case, to build the prototype; scope and protocol agreed in the session.",
     "Confirm willingness in principle; details follow."),
    ("04", "Lock prototype scope within two weeks",
     "A short, precise scope document: use case, data, timeline, and commercial terms. No lengthy proposals.",
     "Target: scope signed by end of week 2."),
]


def render_approach(prs, brand):
    ink = _brand(brand, "ink"); purple = _brand(brand, "primary_purple")
    slate = _brand(brand, "muted"); body = _brand(brand, "body_font")
    slide = _add_blank_slide(prs)
    _header(slide, "Our approach",
            [("A low-risk path to a ", ink), ("live prototype.", purple)], brand,
            subtitle="Value before commitment. Proof before scale.",
            subtitle_color=purple, subtitle_italic=True)
    n = len(_APPROACH_STAGES); cw = 2.22; gap = (11.73 - n * cw) / (n - 1)
    xs = [0.8 + i * (cw + gap) for i in range(n)]; cy = 2.85; ch = 3.5
    for x, (num, when, title, bod) in zip(xs, _APPROACH_STAGES):
        _rrect(slide, x, cy, cw, ch, "F6F5FC", adj=0.04)
        _rect2(slide, x, cy, cw, 0.06, purple)
        _t(slide, x + 0.26, cy + 0.2, cw - 0.4, 0.6, num, font_name=body, font_size=33, bold=True, color_hex=purple)
        _t(slide, x + 0.28, cy + 0.9, cw - 0.4, 0.28, when.upper(), font_name=body, font_size=10, bold=True, color_hex=slate)
        _t(slide, x + 0.28, cy + 1.18, cw - 0.4, 0.42, title, font_name=body, font_size=14.5, bold=True, color_hex=ink)
        _t(slide, x + 0.28, cy + 1.62, cw - 0.45, 1.85, bod, font_name="Calibri Light", font_size=11.5, color_hex=slate)


def render_next_steps(prs, brand, company_name="Company"):
    ink = _brand(brand, "ink"); purple = _brand(brand, "primary_purple")
    slate = _brand(brand, "muted"); body = _brand(brand, "body_font")
    slide = _add_blank_slide(prs)
    _header(slide, "Path forward", [("Path forward.", ink)], brand,
            subtitle="Four things to align on so we can move from this conversation into a working scope.",
            subtitle_color=_brand(brand, "muted"), subtitle_italic=False)
    pos = [(0.8, 2.68), (6.9, 2.68), (0.8, 4.66), (6.9, 4.66)]; iw = 5.63; ih = 1.88
    for (x, y), (num, title, bod, ask) in zip(pos, _NEXT_STEPS):
        _rrect(slide, x, y, iw, ih, "F6F5FC", adj=0.04)
        _rect2(slide, x, y, 0.08, ih, purple)
        _t(slide, x + 0.3, y + 0.16, 0.9, 0.6, num, font_name=body, font_size=27, bold=True, color_hex=purple)
        _t(slide, x + 1.18, y + 0.2, iw - 1.45, 0.4, title, font_name=body, font_size=15.5, bold=True, color_hex=ink)
        _t(slide, x + 1.18, y + 0.62, iw - 1.45, 0.7, bod, font_name="Calibri Light", font_size=12.5, color_hex=slate)
        ay = y + ih - 0.46
        _rect2(slide, x + 1.18, ay - 0.04, 0.42, 0.022, purple)
        _t(slide, x + 1.18, ay, 2.0, 0.2, "THE ASK", font_name=body, font_size=8.5, bold=True, color_hex=purple)
        _t(slide, x + 1.18, ay + 0.19, iw - 1.5, 0.3, ask, font_name="Calibri Light", font_size=11, italic=True, color_hex=slate)


# --------------------------------------------------------------------------------- #
# Delivery capability (role-gap) — heatmap + roles-to-add. Inserted before Our Approach
# when the Talent Gap toggle is on. House chrome via apply_uniform_footer.
# --------------------------------------------------------------------------------- #
_DEP_RAMP = {
    "low":  ("EDE9FE", "6C48F2", "Low"),
    "mod":  ("9B82F0", "FFFFFF", "Mod"),
    "high": ("6C48F2", "FFFFFF", "High"),
}


def render_delivery_capability(prs, content, brand):
    ink = _brand(brand, "ink"); purple = _brand(brand, "primary_purple")
    muted = _brand(brand, "muted"); body = _brand(brand, "body_font")
    rule = _brand(brand, "rule")
    slide = _add_blank_slide(prs)

    # Header
    _pill(slide, 0.8, 0.53, "Delivery capability", purple, "FFFFFF", body)
    _t(slide, 8.0, 0.62, 4.53, 0.3, "ROLE-GAP ASSESSMENT", font_name=body, font_size=10.5,
       bold=True, color_hex=muted, align=PP_ALIGN.RIGHT)
    _two_tone(slide, 0.8, 1.12, 11.9, 0.9,
              [("The roles this transformation ", ink), ("needs.", purple)], body, size=33)
    _t(slide, 0.82, 1.74, 11.6, 0.4,
       "Where the gaps sit today, mapped to the opportunities we'd prioritise.",
       font_name=body, font_size=14.5, italic=True, color_hex=muted)

    body_top = 2.42
    # vertical rule
    vr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.28), Inches(body_top), Pt(1),
                                Inches(5.95 - body_top))
    vr.fill.solid(); vr.fill.fore_color.rgb = _rgb(rule); vr.line.fill.background(); vr.shadow.inherit = False

    roles = content.get("roles", ["FDC", "FDE", "Data Engineer", "FDO"]); cols = content.get("columns", []); hm = content.get("heatmap", [])
    n = max(1, len(cols))

    # LEFT — heatmap
    _t(slide, 0.8, body_top, 5.0, 0.26, "THE TALENT GAP", font_name=body, font_size=11,
       bold=True, color_hex=purple)
    lbl_x = 0.8; lbl_w = 1.45
    grid_x = lbl_x + lbl_w + 0.05
    grid_right = 7.05
    cw = (grid_right - grid_x) / n
    hdr_y = 2.82; hdr_h = 0.46
    row_y0 = hdr_y + hdr_h + 0.04; rh = 0.62

    for j, c in enumerate(cols):
        cx = grid_x + j * cw
        hfs = 9.5
        while hfs > 7 and _est_lines(c, hfs, cw - 0.10) > 2:
            hfs -= 0.5
        _t(slide, cx + 0.02, hdr_y, cw - 0.04, hdr_h, _clamp(c, cw - 0.10, hfs, 2),
           font_name=body, font_size=hfs,
           bold=True, color_hex=ink, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
    for i, role in enumerate(roles):
        ry = row_y0 + i * rh
        _t(slide, lbl_x, ry, lbl_w, rh - 0.10, role, font_name=body, font_size=11,
           bold=True, color_hex=ink, anchor=MSO_ANCHOR.MIDDLE)
        for j in range(n):
            val = hm[i][j] if (i < len(hm) and j < len(hm[i])) else "low"
            bg, tx, lab = _DEP_RAMP.get(val, _DEP_RAMP["low"])
            cx = grid_x + j * cw
            _rrect(slide, cx + 0.03, ry + 0.03, cw - 0.06, rh - 0.12, bg, adj=0.12)
            _t(slide, cx + 0.03, ry + 0.02, cw - 0.06, rh - 0.12, lab, font_name=body,
               font_size=10, bold=True, color_hex=tx, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # RIGHT — roles to add
    rx = 7.5; rw = 5.03
    _t(slide, rx, body_top, rw, 0.26, "ROLES TO ADD", font_name=body, font_size=11,
       bold=True, color_hex=purple)
    recs = content.get("recommendations", [])[:3]
    cy = 2.82; ch = 1.15; gap = 0.10
    for rec in recs:
        _rrect(slide, rx, cy, rw, ch, "FAFAFC", adj=0.05, line="ECECF3", line_w=1.0)
        _rect2(slide, rx, cy, 0.06, ch, purple)
        _t(slide, rx + 0.25, cy + 0.11, rw - 0.45, 0.3, rec.get("role", ""), font_name=body,
           font_size=12.5, bold=True, color_hex=ink)
        _t(slide, rx + 0.25, cy + 0.40, rw - 0.45, 0.34,
           _clamp(rec.get("contribution", ""), rw - 0.45, 10, 2),
           font_name="Calibri Light", font_size=10, color_hex=muted)
        hy = cy + 0.84
        _t(slide, rx + 0.25, hy, 0.5, 0.2, "HELPS", font_name=body, font_size=8,
           bold=True, color_hex=purple)
        chx = rx + 0.74
        for tag in rec.get("helps", []):
            tw = len(tag) * 0.066 + 0.22
            if chx + tw > rx + rw - 0.15:   # would overflow the card — stop here
                break
            chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(chx), Inches(hy - 0.035),
                                          Inches(tw), Inches(0.24))
            chip.fill.solid(); chip.fill.fore_color.rgb = _rgb("F0F0F4")
            chip.line.fill.background(); chip.shadow.inherit = False
            try:
                chip.adjustments[0] = 0.5
            except Exception:
                pass
            _t(slide, chx, hy - 0.04, tw, 0.24, tag, font_name=body, font_size=8.5,
               color_hex=muted, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            chx += tw + 0.1
        cy += ch + gap

    # method note (sources style)
    _t(slide, 0.8, 6.5, 10.5, 0.3, content.get("method_note", ""), font_name="Arial",
       font_size=10.0, italic=True, color_hex=muted)
